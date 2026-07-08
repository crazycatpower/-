# -*- coding: utf-8 -*-
import argparse
import io
import json
import os
import re
import shutil
import socket
import subprocess
import tempfile
import threading
from datetime import datetime, timedelta
from pathlib import Path

from dotenv import load_dotenv
# 必須在匯入 db / safety_audit_system 等本地模組之前載入 .env，
# 否則這些模組在頂層讀取的 os.getenv(...) 會拿到空值。
load_dotenv(dotenv_path=str(Path(__file__).resolve().parent / ".env"), override=True)

from flask import Flask, jsonify, request, send_from_directory, send_file

# LINE SDK v3
from linebot.v3.webhook import WebhookHandler
from linebot.v3.exceptions import InvalidSignatureError
from linebot.v3.webhooks import MessageEvent, TextMessageContent, ImageMessageContent, FileMessageContent
from linebot.v3.messaging import (
    ApiClient, Configuration, MessagingApi, MessagingApiBlob,
    ReplyMessageRequest, PushMessageRequest,
    TextMessage, ImageMessage,
)
try:
    from linebot.v3.messaging import (
        FlexMessage, FlexBubble, FlexBox, FlexText, FlexSeparator,
        FlexButton, FlexCarousel,
        QuickReply, QuickReplyItem, MessageAction,
    )
    _FLEX_OK = True
except ImportError:
    _FLEX_OK = False

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as _plt
    _MATPLOTLIB_OK = True
except ImportError:
    _MATPLOTLIB_OK = False

# AI 分析模組
try:
    from safety_audit_system import analyze_image
except ImportError:
    def analyze_image(**kwargs):
        return {"scene_analysis": {"summary": "分析模組異常"}, "detections": [], "annotated_image": ""}

from audit_checklist import official_text

# 法規知識庫（選填）
try:
    from knowledge_base import retrieve_regulations, list_regulation_sources
    KNOWLEDGE_BASE_AVAILABLE = True
except ImportError:
    KNOWLEDGE_BASE_AVAILABLE = False
    def list_regulation_sources():
        return []

# 資料庫（選填）
try:
    from db import (
        DB_RUN_BACKFILL, DB_RUN_MIGRATIONS, build_where,
        db_conn, release_advisory_lock, with_advisory_lock,
    )
    from psycopg.sql import SQL
    from psycopg.types.json import Jsonb
    DB_AVAILABLE = True
except ImportError:
    DB_AVAILABLE = False
    DB_RUN_MIGRATIONS = False
    DB_RUN_BACKFILL = False

# --- 路徑設定 ---
BASE_DIR = Path(__file__).resolve().parent
CERT_DIR = BASE_DIR / "certs"
CERT_FILE = CERT_DIR / "dev-cert.pem"
KEY_FILE = CERT_DIR / "dev-key.pem"
REGULATIONS_DIR = BASE_DIR / "regulations"

CERT_DIR.mkdir(exist_ok=True)

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 10 * 1024 * 1024

# --- LINE 憑證（從 .env 讀取） ---
LINE_CHANNEL_ACCESS_TOKEN = (os.getenv("CHANNEL_ACCESS_TOKEN") or "").strip()
LINE_CHANNEL_SECRET = (os.getenv("CHANNEL_SECRET") or "").strip()
if not LINE_CHANNEL_ACCESS_TOKEN or not LINE_CHANNEL_SECRET:
    raise RuntimeError("請在 .env 設定 CHANNEL_ACCESS_TOKEN 與 CHANNEL_SECRET")
PUBLIC_URL = (os.getenv("PUBLIC_URL") or "").strip()

_CLOUDFLARED_DEFAULT = shutil.which("cloudflared") or r"C:\Program Files (x86)\cloudflared\cloudflared.exe"
CLOUDFLARED_PATH = (os.getenv("CLOUDFLARED_PATH") or _CLOUDFLARED_DEFAULT).strip()
_cloudflared_proc: subprocess.Popen | None = None

_FALLBACK_PROMPT = "請分析目前畫面，重點說明可能的標誌、風險與建議行動。"


def looks_like_mojibake(text: str) -> bool:
    suspicious = ("é¢", "å", "ç", "é", "è", "î¼¾")
    value = str(text or "").strip()
    return not value or any(t in value for t in suspicious)


DEFAULT_ANALYZE_PROMPT = (os.getenv("DEFAULT_ANALYZE_PROMPT") or _FALLBACK_PROMPT).strip()
if looks_like_mojibake(DEFAULT_ANALYZE_PROMPT):
    DEFAULT_ANALYZE_PROMPT = _FALLBACK_PROMPT

_line_config = Configuration(access_token=LINE_CHANNEL_ACCESS_TOKEN)
handler = WebhookHandler(LINE_CHANNEL_SECRET)


def _line_reply(reply_token: str, *messages) -> None:
    with ApiClient(_line_config) as api_client:
        MessagingApi(api_client).reply_message(
            ReplyMessageRequest(reply_token=reply_token, messages=list(messages))
        )


def _line_push(to: str, *messages) -> None:
    with ApiClient(_line_config) as api_client:
        MessagingApi(api_client).push_message(
            PushMessageRequest(to=to, messages=list(messages))
        )


def _qr(*items: tuple) -> QuickReply | None:
    """Build a QuickReply from (label, text) pairs. Returns None if Flex unavailable."""
    if not _FLEX_OK:
        return None
    return QuickReply(items=[
        QuickReplyItem(action=MessageAction(label=str(label), text=str(text)))
        for label, text in items[:13]
    ])


def _line_get_content(message_id: str) -> bytes:
    with ApiClient(_line_config) as api_client:
        return bytes(MessagingApiBlob(api_client).get_message_content(message_id))

UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

# --- 六種模式 ---
MODE_IMAGE_ANALYSIS   = "image_analysis"
MODE_REGULATION_QUERY = "regulation_query"
MODE_CASE_STUDY       = "case_study"
MODE_AUDIT_RECORD     = "audit_record"
MODE_KNOWLEDGE_BASE   = "knowledge_base"

LINE_USER_STATE: dict[str, dict[str, str]] = {}

MODE_COMMANDS: dict[str, dict] = {
    "法規查詢": {"mode": MODE_REGULATION_QUERY},
    "案例學習": {"mode": MODE_CASE_STUDY},
    "稽核紀錄": {"mode": MODE_AUDIT_RECORD},
    "照片稽核": {"mode": MODE_IMAGE_ANALYSIS, "prompt": DEFAULT_ANALYZE_PROMPT},
    "知識庫":   {"mode": MODE_KNOWLEDGE_BASE},
    "風險分析": {"mode": MODE_IMAGE_ANALYSIS, "prompt": "請分析這張照片中的風險、異常與建議處置方式。"},
    "標誌辨識": {"mode": MODE_IMAGE_ANALYSIS, "prompt": "請辨識畫面中的交通標誌或安全標示，說明其意義與注意事項。"},
}
HELP_COMMANDS = {"help", "menu", "幫助", "說明"}
NEXT_PAGE_COMMANDS = {"下一頁", "next"}
PREV_PAGE_COMMANDS = {"上一頁", "prev"}
PAGE_SIZE = 5

_MODE_LABELS = {
    MODE_IMAGE_ANALYSIS:   "照片稽核",
    MODE_REGULATION_QUERY: "法規查詢",
    MODE_CASE_STUDY:       "案例學習",
    MODE_AUDIT_RECORD:     "稽核紀錄",
    MODE_KNOWLEDGE_BASE:   "知識庫",
}
_MODE_HINTS = {
    MODE_REGULATION_QUERY: "請直接輸入法規問題，支援翻頁。",
    MODE_IMAGE_ANALYSIS:   "請傳送照片，系統會自動進行 AI 分析。",
    MODE_AUDIT_RECORD:     "輸入時間查詢紀錄（如 2026-03-31），或輸入「統計」查看整體數據。",
    MODE_CASE_STUDY:       "輸入「概覽」查最近案例分析，「本週」/「本月」查趨勢，「最常見」查常見違規。",
    MODE_KNOWLEDGE_BASE:   "輸入「新增」後傳送檔案（PDF/Word/Excel…）新增法規；「列出」查看來源；「刪除 [來源]」刪除；「匯入」從資料夾批次匯入。",
}
_MODE_IMAGE_REJECTS = {
    MODE_REGULATION_QUERY: "法規查詢模式請使用文字提問。",
    MODE_CASE_STUDY:       "案例學習模式請輸入文字指令。",
    MODE_AUDIT_RECORD:     "稽核紀錄模式請直接輸入時間查詢。",
    MODE_KNOWLEDGE_BASE:   "知識庫模式請輸入文字指令。",
}


# =============================================================================
# 使用者狀態管理
# =============================================================================

def get_source_key(event) -> str:
    source = event.source
    for attr in ("user_id", "group_id", "room_id"):
        value = getattr(source, attr, None)
        if value:
            return f"{source.type}:{value}"
    return str(source.type)


def get_user_state(source_key: str) -> dict:
    if source_key not in LINE_USER_STATE:
        LINE_USER_STATE[source_key] = {
            "mode": MODE_IMAGE_ANALYSIS,
            "prompt": DEFAULT_ANALYZE_PROMPT,
            "page_items": [],
            "page_index": 0,
            "page_query": "",
            "page_mode": "",
            "kb_pending": None,  # {"action": "add", "source": "name"} when waiting for content
        }
    return LINE_USER_STATE[source_key]


def set_user_mode(source_key: str, mode: str, prompt: str | None = None) -> None:
    state = get_user_state(source_key)
    state["mode"] = mode
    if prompt is not None:
        state["prompt"] = prompt


# =============================================================================
# 通用工具函式
# =============================================================================

def read_file_bytes(path_like) -> bytes:
    if not path_like:
        return b""
    try:
        return Path(path_like).read_bytes()
    except OSError:
        return b""


# =============================================================================
# 文件文字擷取（知識庫上傳用）
# =============================================================================

_SUPPORTED_EXTS = {
    ".txt", ".md", ".csv",
    ".pdf",
    ".docx", ".doc",
    ".xlsx", ".xls",
    ".pptx",
    ".html", ".htm",
}


def _decode_text(data: bytes) -> str:
    try:
        import chardet
        enc = (chardet.detect(data) or {}).get("encoding") or "utf-8"
    except ImportError:
        enc = "utf-8"
    return data.decode(enc, errors="replace")


def _extract_pdf(data: bytes) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        return "\n\n".join(p.extract_text() or "" for p in reader.pages)
    except ImportError:
        raise RuntimeError("請安裝 pypdf：pip install pypdf")


def _extract_docx(data: bytes) -> str:
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        raise RuntimeError("請安裝 python-docx：pip install python-docx")


def _extract_xlsx(data: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        parts = []
        for name in wb.sheetnames:
            rows = []
            for row in wb[name].iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells))
            if rows:
                parts.append(f"【{name}】\n" + "\n".join(rows))
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("請安裝 openpyxl：pip install openpyxl")


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"【第 {i} 頁】\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except ImportError:
        raise RuntimeError("請安裝 python-pptx：pip install python-pptx")


def _extract_html(data: bytes) -> str:
    try:
        from bs4 import BeautifulSoup
        return BeautifulSoup(data, "html.parser").get_text(separator="\n", strip=True)
    except ImportError:
        text = _decode_text(data)
        return re.sub(r"<[^>]+>", " ", text)


def extract_text_from_file(data: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext not in _SUPPORTED_EXTS:
        raise RuntimeError(
            f"不支援的檔案格式「{ext or '無副檔名'}」。\n"
            f"支援格式：{', '.join(sorted(_SUPPORTED_EXTS))}"
        )
    if ext in (".txt", ".md", ".csv"):
        return _decode_text(data)
    if ext == ".pdf":
        return _extract_pdf(data)
    if ext in (".docx", ".doc"):
        return _extract_docx(data)
    if ext in (".xlsx", ".xls"):
        return _extract_xlsx(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    if ext in (".html", ".htm"):
        return _extract_html(data)
    return _decode_text(data)


def get_report_path_from_result(result: dict) -> str:
    annotated = result.get("annotated_image")
    if not annotated:
        return ""
    try:
        p = Path(annotated).resolve().with_name("report.json")
        if OUTPUT_DIR.resolve() not in p.parents:
            return ""
        return str(p)
    except Exception:
        return ""


def load_report_json(result: dict) -> dict:
    embedded = result.get("report_json")
    if isinstance(embedded, dict):
        return embedded
    report_path = get_report_path_from_result(result)
    if not report_path:
        return {}
    try:
        data = json.loads(Path(report_path).read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return data if isinstance(data, dict) else {"data": data}


_CATEGORY_ICONS = {
    "墜落防止":    "🪜",
    "倒塌崩塌防止": "🏗",
    "感電防止":    "⚡",
    "火災爆炸":    "🔥",
    "中毒缺氧":    "☠️",
    "交通維持":    "🚦",
    "其他":       "📋",
}
_RESULT_LABELS  = {"○": "符合", "×": "缺失", "△": "待改善"}
_MARKER_LABELS  = {"☆": "【須停工】", "※": "【即日改善】", "": "【五日內改善】"}
_RESULT_TICK    = {"○": "✅", "×": "❌", "△": "⚠️"}

# 這6類會逐條打勾完整列出（對應 audit_checklist.CHECKLIST_1_TO_6），
# 其餘（目前只有「其他」）維持只列缺失，不逐條窮舉。
_FULL_CHECKLIST_CATEGORIES = {"墜落防止", "倒塌崩塌防止", "感電防止", "火災爆炸", "中毒缺氧", "交通維持"}


def _format_categories_text(scene: dict, categories: dict) -> str:
    parts: list[str] = []
    summary = str(scene.get("summary") or "").strip()
    if summary:
        parts.append(f"📝 {summary}")

    has_violations = False
    for cat_name, items in categories.items():
        if not items:
            continue
        icon = _CATEGORY_ICONS.get(cat_name, "📋")

        if cat_name in _FULL_CHECKLIST_CATEGORIES:
            # 逐條打勾清單：符合的項目只顯示一行勾選，缺失才展開細節
            n_defects = sum(1 for i in items if isinstance(i, dict) and i.get("result") in ("×", "△"))
            parts.append(f"\n{icon}【{cat_name}】（{len(items)}項，缺失{n_defects}項）")
            for it in items:
                if not isinstance(it, dict):
                    continue
                result  = it.get("result", "○")
                item_id = it.get("item", "")
                tick    = _RESULT_TICK.get(result, "▫️")
                if result == "○":
                    parts.append(f"  {tick} {item_id}")
                    continue
                has_violations = True
                marker   = it.get("marker", "")
                desc     = str(it.get("description", "")).strip()
                location = str(it.get("location", "")).strip()
                deadline = str(it.get("deadline", "")).strip()
                rule_text = official_text(item_id)
                r_label  = _RESULT_LABELS.get(result, result)
                m_label  = _MARKER_LABELS.get(marker, "")
                line     = f"  {tick} {item_id} [{r_label}]{m_label}"
                if rule_text:
                    line += f"\n     📖 {rule_text}"
                if desc:
                    line += f"\n     {desc}"
                if location:
                    line += f"\n     📍 地點：{location}"
                if deadline:
                    line += f"\n     ⏰ 期限：{deadline}"
                parts.append(line)
            continue

        # 「其他」等非逐條類別：維持只列有缺失的項目
        violations = [i for i in items if isinstance(i, dict) and i.get("result") in ("×", "△")]
        if not violations:
            continue
        has_violations = True
        parts.append(f"\n{icon}【{cat_name}】")
        for it in violations:
            result   = it.get("result", "")
            marker   = it.get("marker", "")
            item_id  = it.get("item", "")
            desc     = str(it.get("description", "")).strip()
            location = str(it.get("location", "")).strip()
            deadline = str(it.get("deadline", "")).strip()
            rule_text = official_text(item_id)
            r_label  = _RESULT_LABELS.get(result, result)
            m_label  = _MARKER_LABELS.get(marker, "")
            line     = f"  {marker}{item_id} [{r_label}]{m_label}"
            if rule_text:
                line += f"\n    📖 {rule_text}"
            if desc:
                line += f"\n    {desc}"
            if location:
                line += f"\n    📍 地點：{location}"
            if deadline:
                line += f"\n    ⏰ 期限：{deadline}"
            parts.append(line)

    if not has_violations:
        parts.append("\n✅ 本次稽核未發現明顯缺失")

    next_action = str(scene.get("next_action") or "").strip()
    if next_action:
        parts.append(f"\n⚖️ 改善建議：{next_action}")

    refs = scene.get("regulation_refs") or []
    if refs:
        flat = [
            str(r.get("cite", r.get("source", str(r)))) if isinstance(r, dict) else str(r)
            for r in refs[:4]
        ]
        parts.append(f"📚 法規依據：{'；'.join(flat)[:300]}")

    return "\n".join(parts)


def scene_audit_json_text(scene: dict) -> str:
    if not scene:
        return ""
    categories = scene.get("categories")
    if isinstance(categories, dict):
        return _format_categories_text(scene, categories)
    # 舊格式相容
    parts: list[str] = []
    summary = str(scene.get("summary") or "").strip()
    if summary:
        parts.append(f"📝 摘要：{summary}")
    risks = [str(r).strip() for r in (scene.get("risks") or []) if str(r).strip()]
    if risks:
        parts.append("⚠️ 風險：" + "、".join(risks[:5]))
    next_action = str(scene.get("next_action") or "").strip()
    if next_action:
        parts.append(f"⚖️ 建議：{next_action}")
    refs = scene.get("regulation_refs") or []
    if refs:
        flat: list[str] = []
        for r in refs[:4]:
            if isinstance(r, dict):
                flat.append(str(r.get("cite", r.get("source", str(r)))))
            else:
                flat.append(str(r))
        parts.append("📚 法規：" + "；".join(flat)[:300])
    answer = str(scene.get("answer") or "").strip()
    if answer:
        parts.append(f"💡 {answer[:300]}")
    return "\n".join(parts)


# =============================================================================
# 圖片分析
# =============================================================================

def analyze_image_bytes(
    image_bytes: bytes,
    *,
    prompt: str,
    use_yolo: bool = True,
    filename_stem: str = "capture",
    extension: str = ".jpg",
) -> dict:
    with tempfile.NamedTemporaryFile(
        suffix=extension, prefix=f"{filename_stem}_", dir=UPLOAD_DIR, delete=False
    ) as tmp:
        tmp.write(image_bytes)
        tmp_path = Path(tmp.name)
    try:
        model_path = BASE_DIR / "best.pt"
        if not model_path.exists():
            model_path = BASE_DIR / "yolov8n.pt"
        result = analyze_image(
            image_path=tmp_path,
            yolo_model_path=str(model_path),
            output_dir=OUTPUT_DIR,
            prompt=prompt or DEFAULT_ANALYZE_PROMPT,
            use_yolo=use_yolo,
        )
    finally:
        tmp_path.unlink(missing_ok=True)
    result["annotated_image_bytes"] = read_file_bytes(result.get("annotated_image"))
    result["report_json"] = load_report_json(result)
    result["input_image_bytes"] = image_bytes
    for _p in [result.get("annotated_image"), get_report_path_from_result(result)]:
        if _p:
            Path(_p).unlink(missing_ok=True)
    return result


def summarize_scene_analysis(result: dict) -> str:
    scene = result.get("scene_analysis") or {}
    compact = scene_audit_json_text(scene)
    if compact:
        return compact[:1500]
    parts: list[str] = []
    summary = str(scene.get("summary") or scene.get("answer") or "").strip()
    if summary:
        parts.append(summary)
    risks = [str(r).strip() for r in (scene.get("risks") or []) if str(r).strip()]
    if risks:
        parts.append("風險：" + "、".join(risks[:3]))
    next_action = str(scene.get("next_action") or "").strip()
    if next_action:
        parts.append("建議：" + next_action)
    return "\n".join(parts)[:1500]


def _summarize_detections(detections: list[dict]) -> str:
    """把 YOLO 偵測清單濃縮成「類別 x數量」，取代單純的物件總數字，
    讓使用者一眼看出偵測到的是什麼（如「安全錐 x4」），而不只是「4 項目」。"""
    if not detections:
        return "無"
    from collections import Counter
    counts = Counter(str(d.get("yolo_class") or d.get("class") or "未知") for d in detections)
    return "、".join(f"{cls} x{n}" for cls, n in counts.most_common())


def build_line_analysis_reply(result: dict, prompt: str) -> str:
    scene = result.get("scene_analysis") or {}
    risk_level, _, risk_emoji = assess_risk_level(result)
    detections = result.get("detections") or []
    header = (
        f"{risk_emoji} 風險等級：{risk_level}　"
        f"偵測到 {len(detections)} 項目（{_summarize_detections(detections)}）"
    )
    compact = scene_audit_json_text(scene)
    if compact:
        return f"{header}\n\n{compact}"[:4900]
    lines = [header]
    if prompt:
        lines.append(f"分析需求：{prompt}")
    summary = str(scene.get("summary") or scene.get("answer") or "").strip()
    if summary:
        lines.append(f"摘要：{summary}")
    risks = [str(r).strip() for r in (scene.get("risks") or []) if str(r).strip()]
    if risks:
        lines.append("風險：" + "、".join(risks[:3]))
    next_action = str(scene.get("next_action") or "").strip()
    if next_action:
        lines.append(f"建議：{next_action}")
    return "\n".join(lines)[:4900]


# =============================================================================
# 資料庫
# =============================================================================

def normalize_audit_time_query(text: str) -> str:
    return "".join(ch for ch in str(text or "").strip() if ch.isdigit())


def init_audit_db() -> None:
    if not DB_AVAILABLE or not DB_RUN_MIGRATIONS:
        return
    with db_conn() as conn:
        with_advisory_lock(conn, "audit_records_migrations_v1")
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    CREATE TABLE IF NOT EXISTS audit_records (
                        id BIGSERIAL PRIMARY KEY,
                        created_at TEXT NOT NULL,
                        source TEXT NOT NULL,
                        source_key TEXT NOT NULL,
                        analysis_mode TEXT NOT NULL,
                        prompt TEXT NOT NULL,
                        image_path TEXT NOT NULL DEFAULT '',
                        annotated_image TEXT NOT NULL DEFAULT '',
                        report_path TEXT NOT NULL DEFAULT '',
                        created_at_digits TEXT NOT NULL DEFAULT '',
                        detection_count INTEGER NOT NULL DEFAULT 0,
                        summary TEXT NOT NULL DEFAULT '',
                        original_image_data BYTEA,
                        reference_image_data BYTEA,
                        annotated_image_data BYTEA,
                        report_json JSONB NOT NULL DEFAULT '{}'::jsonb
                    )
                """)
                for stmt in [
                    "ALTER TABLE audit_records ADD COLUMN IF NOT EXISTS created_at_digits TEXT NOT NULL DEFAULT ''",
                    "ALTER TABLE audit_records ADD COLUMN IF NOT EXISTS original_image_data BYTEA",
                    "ALTER TABLE audit_records ADD COLUMN IF NOT EXISTS annotated_image_data BYTEA",
                    "ALTER TABLE audit_records ADD COLUMN IF NOT EXISTS report_json JSONB NOT NULL DEFAULT '{}'::jsonb",
                    "ALTER TABLE audit_records ADD COLUMN IF NOT EXISTS reference_image_data BYTEA",
                    "UPDATE audit_records SET created_at_digits = regexp_replace(created_at, '[^0-9]', '', 'g') WHERE created_at_digits = ''",
                    "UPDATE audit_records SET report_json = '{}'::jsonb WHERE report_json IS NULL",
                    "CREATE INDEX IF NOT EXISTS idx_audit_records_created_at ON audit_records (created_at)",
                    "CREATE INDEX IF NOT EXISTS idx_audit_records_source_key ON audit_records (source_key)",
                    "CREATE INDEX IF NOT EXISTS idx_audit_records_created_at_digits ON audit_records (created_at_digits)",
                    "CREATE INDEX IF NOT EXISTS idx_audit_records_source_key_created_at ON audit_records (source_key, created_at DESC)",
                ]:
                    cur.execute(stmt)
            conn.commit()
            if DB_RUN_BACKFILL:
                backfill_audit_record_binary_data(conn)
        finally:
            release_advisory_lock(conn, "audit_records_migrations_v1")


def save_audit_record(
    *,
    source: str,
    source_key: str,
    prompt: str,
    original_image_data: bytes,
    result: dict,
    analysis_mode: str,
    reference_image_data: bytes | None = None,
) -> dict:
    created_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    report_json = load_report_json(result)
    record: dict = {
        "created_at": created_at,
        "source": source,
        "source_key": source_key,
        "analysis_mode": analysis_mode,
        "prompt": prompt,
        "original_image_data": original_image_data,
        "reference_image_data": reference_image_data,
        "annotated_image_data": bytes(result.get("annotated_image_bytes") or b""),
        "report_json": report_json,
        "created_at_digits": normalize_audit_time_query(created_at),
        "detection_count": len(result.get("detections") or []),
        "summary": summarize_scene_analysis(result),
        "id": None,
    }
    if not DB_AVAILABLE:
        return record
    with db_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO audit_records (
                    created_at, source, source_key, analysis_mode, prompt,
                    image_path, annotated_image, report_path, created_at_digits,
                    detection_count, summary, original_image_data,
                    annotated_image_data, report_json, reference_image_data
                ) VALUES (%s, %s, %s, %s, %s, '', '', '', %s, %s, %s, %s, %s, %s, %s)
                RETURNING id
                """,
                (
                    record["created_at"], record["source"], record["source_key"],
                    record["analysis_mode"], record["prompt"],
                    record["created_at_digits"], record["detection_count"],
                    record["summary"], record["original_image_data"],
                    record["annotated_image_data"], Jsonb(report_json),
                    record["reference_image_data"],
                ),
            )
            inserted = cur.fetchone()
            record["id"] = inserted["id"] if inserted else None
        conn.commit()
    return record


def find_audit_records(
    time_query: str, source_key: str | None = None, limit: int = 10, offset: int = 0
) -> list[dict]:
    if not DB_AVAILABLE:
        return []
    normalized = normalize_audit_time_query(time_query)
    with db_conn() as conn:
        where_clauses: list = ["TRUE"]
        params: list = []
        if source_key:
            where_clauses.append("source_key IN ('', %s)")
            params.append(source_key)
        if normalized:
            where_clauses.append("created_at_digits LIKE %s")
            params.append(f"%{normalized}%")
        params.extend([limit, offset])
        query = SQL("""
            SELECT id, created_at, source, source_key, analysis_mode, prompt,
                   detection_count, summary, report_json,
                   OCTET_LENGTH(original_image_data) > 0 AS has_original_image,
                   OCTET_LENGTH(reference_image_data) > 0 AS has_reference_image,
                   OCTET_LENGTH(annotated_image_data) > 0 AS has_annotated_image
            FROM audit_records
            WHERE {where}
            ORDER BY created_at DESC
            LIMIT %s OFFSET %s
        """).format(where=build_where(where_clauses))
        with conn.cursor() as cur:
            cur.execute(query, params)
            return cur.fetchall()


def get_case_study_stats(days: int = 30) -> dict:
    if not DB_AVAILABLE:
        return {}
    since = (datetime.now() - timedelta(days=days)).strftime("%Y-%m-%d")
    keywords = ["安全帽", "安全帶", "鷹架", "護欄", "反光背心", "安全鞋", "電氣", "火災", "墜落", "感電"]
    with db_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT COUNT(*) AS total,
                       ROUND(AVG(detection_count)::numeric, 1) AS avg_det,
                       MAX(detection_count) AS max_det,
                       SUM(detection_count) AS total_det
                FROM audit_records WHERE SUBSTR(created_at, 1, 10) >= %s
                """,
                (since,),
            )
            overview = cur.fetchone() or {}
            # 各關鍵字命中次數
            cur.execute(
                f"SELECT {', '.join(f'SUM(CASE WHEN summary ILIKE %s THEN 1 ELSE 0 END) AS k{i}' for i, _ in enumerate(keywords))} "
                f"FROM audit_records WHERE SUBSTR(created_at, 1, 10) >= %s",
                [f"%{k}%" for k in keywords] + [since],
            )
            kw_row = cur.fetchone() or {}
            kw_counts = sorted(
                [(keywords[i], kw_row.get(f"k{i}") or 0) for i in range(len(keywords))],
                key=lambda x: x[1], reverse=True,
            )
            # 每日趨勢
            cur.execute(
                """
                SELECT SUBSTR(created_at, 1, 10) AS day,
                       COUNT(*) AS cnt, SUM(detection_count) AS detections
                FROM audit_records WHERE SUBSTR(created_at, 1, 10) >= %s
                GROUP BY day ORDER BY day DESC LIMIT 14
                """,
                (since,),
            )
            trend = cur.fetchall()
    return {"overview": overview, "kw_counts": kw_counts, "trend": trend, "since": since}


def backfill_audit_record_binary_data(conn=None) -> None:
    if not DB_AVAILABLE:
        return
    if conn is None:
        with db_conn() as c:
            return backfill_audit_record_binary_data(c)
    with conn.cursor() as cur:
        cur.execute("""
            SELECT id, image_path, annotated_image, report_path,
                   original_image_data, annotated_image_data, reference_image_data, report_json
            FROM audit_records
            WHERE (original_image_data IS NULL AND image_path <> '')
               OR (annotated_image_data IS NULL AND annotated_image <> '')
               OR (report_json = '{}'::jsonb AND report_path <> '')
        """)
        rows = cur.fetchall()
        for row in rows:
            orig = row.get("original_image_data")
            ann = row.get("annotated_image_data")
            rj = row.get("report_json") or {}
            if orig is None and row.get("image_path"):
                orig = read_file_bytes(row["image_path"])
            if ann is None and row.get("annotated_image"):
                ann = read_file_bytes(row["annotated_image"])
            if rj == {} and row.get("report_path"):
                try:
                    loaded = json.loads(Path(row["report_path"]).read_text(encoding="utf-8"))
                    rj = loaded if isinstance(loaded, dict) else {"data": loaded}
                except (OSError, json.JSONDecodeError):
                    rj = {}
            cur.execute(
                """
                UPDATE audit_records
                SET original_image_data = COALESCE(%s, original_image_data),
                    annotated_image_data = COALESCE(%s, annotated_image_data),
                    report_json = CASE WHEN report_json = '{}'::jsonb THEN %s ELSE report_json END
                WHERE id = %s
                """,
                (orig, ann, Jsonb(rj), row["id"]),
            )
    conn.commit()


# =============================================================================
# 風險等級評估
# =============================================================================

def assess_risk_level(result: dict) -> tuple[str, str, str]:
    """回傳 (等級文字, Flex 標題背景色, emoji)。

    只依 categories 裡「真的有缺失」的項目（result 為 × 或 △）之 marker 判斷等級，
    不再用關鍵字比對 summary/answer 全文。categories 現在是逐條打勾的完整清單，
    這些文字欄位本來就會提到「墜落」「感電」「火災」等類別名稱——即使該類別完全
    符合規定、答案是在說「現場符合墜落、感電防護規定」，關鍵字比對一樣會命中，
    完全無法反映有沒有真的缺失，實測幾乎每次都被誤判成「高風險」。
    """
    scene = result.get("scene_analysis") or {}
    categories = scene.get("categories") or {}

    defect_markers = {
        str(item.get("marker") or "")
        for items in categories.values()
        for item in (items or [])
        if isinstance(item, dict) and item.get("result") in ("×", "△")
    }

    if "☆" in defect_markers:
        return "高", "#991b1b", "🔴"
    if "※" in defect_markers:
        return "中", "#92400e", "🟠"
    if defect_markers:
        # 空字串 marker（五日內改善）的缺失，或舊格式沒有 marker 的一般缺失
        return "低", "#14532d", "🟡"

    det_count = len(result.get("detections") or [])
    if det_count >= 1:
        return "低", "#14532d", "🟡"
    return "低", "#14532d", "🟢"


# =============================================================================
# LINE Flex Message 卡片
# =============================================================================

def _flex_category_contents(categories: dict) -> list:
    """把 categories dict 轉成 Flex body contents。

    「墜落防止~交通維持」這6類逐條打勾列出（符合的項目濃縮成一個勾選 chip，
    只有缺失才展開細節），呈現成像官方稽核表一樣「一條一條看」的完整清單；
    「其他」等非逐條類別維持只列缺失。
    """
    _RESULT_COLORS = {"×": "#dc2626", "△": "#d97706", "○": "#16a34a"}
    _RESULT_BADGE = {"×": "✗", "△": "△", "○": "✓"}
    contents = []
    has_any = False

    for cat_name, items in categories.items():
        items = items or []
        if not items:
            continue
        icon = _CATEGORY_ICONS.get(cat_name, "📋")

        if cat_name in _FULL_CHECKLIST_CATEGORIES:
            n_defects = sum(1 for i in items if isinstance(i, dict) and i.get("result") in ("×", "△"))
            contents.append(FlexSeparator(margin="md"))
            contents.append(FlexText(
                text=f"{icon} {cat_name}（缺失 {n_defects}/{len(items)}）",
                weight="bold", size="sm", color="#1e3a5f", margin="md",
            ))
            # 符合規定的項目濃縮成一行逗號分隔的勾選清單，避免卡片被 54 條全展開撐爆
            compliant_ids = [str(i.get("item", "")) for i in items if isinstance(i, dict) and i.get("result") == "○"]
            if compliant_ids:
                contents.append(FlexText(
                    text="✓ 符合：" + "、".join(compliant_ids),
                    size="xs", color="#16a34a", wrap=True,
                ))
            for it in items:
                if not isinstance(it, dict) or it.get("result") not in ("×", "△"):
                    continue
                has_any = True
                result  = it.get("result", "")
                marker  = it.get("marker", "")
                item_id = it.get("item", "")
                desc    = str(it.get("description", "")).strip()
                location = str(it.get("location", "")).strip()
                deadline = str(it.get("deadline", "")).strip()
                r_color  = _RESULT_COLORS.get(result, "#374151")
                badge    = _RESULT_BADGE.get(result, result)
                label    = f"{marker}{item_id} [{badge}]"
                if location:
                    label += f"　📍{location}"
                contents.append(FlexText(text=label, size="xs", color=r_color, weight="bold"))
                if desc:
                    contents.append(FlexText(text=desc[:120], wrap=True, size="xs", color="#374151"))
                if deadline:
                    contents.append(FlexText(text=f"⏰ {deadline}", size="xs", color="#6b7280"))
            continue

        violations = [i for i in items if isinstance(i, dict) and i.get("result") in ("×", "△")]
        if not violations:
            continue
        has_any = True
        contents.append(FlexSeparator(margin="md"))
        contents.append(FlexText(
            text=f"{icon} {cat_name}",
            weight="bold", size="sm", color="#1e3a5f", margin="md",
        ))
        for it in violations[:4]:
            result  = it.get("result", "")
            marker  = it.get("marker", "")
            item_id = it.get("item", "")
            desc    = str(it.get("description", "")).strip()
            location = str(it.get("location", "")).strip()
            deadline = str(it.get("deadline", "")).strip()
            r_color  = _RESULT_COLORS.get(result, "#374151")
            badge    = _RESULT_BADGE.get(result, result)
            label    = f"{marker}{item_id} [{badge}]"
            if location:
                label += f"　📍{location}"
            contents.append(FlexText(text=label, size="xs", color=r_color, weight="bold"))
            if desc:
                contents.append(FlexText(text=desc[:120], wrap=True, size="xs", color="#374151"))
            if deadline:
                contents.append(FlexText(text=f"⏰ {deadline}", size="xs", color="#6b7280"))

    if not has_any:
        contents.append(FlexText(text="✅ 本次稽核未發現明顯缺失", size="sm", color="#16a34a"))
    return contents


def build_flex_analysis_message(result: dict, audit_id: int | None = None):
    """把圖片分析結果包成 Flex Message 氣泡卡片（稽核表格式）；失敗時回傳 None。"""
    if not _FLEX_OK:
        return None
    try:
        scene = result.get("scene_analysis") or {}
        det_count = len(result.get("detections") or [])
        risk_level, header_color, risk_emoji = assess_risk_level(result)
        body_contents = []

        body_contents.append(FlexText(
            text=f"偵測到 {det_count} 個項目　風險等級：{risk_emoji} {risk_level}",
            size="sm", color="#6b7280",
        ))

        summary = str(scene.get("summary") or scene.get("answer") or "").strip()
        if summary:
            body_contents.append(FlexSeparator(margin="md"))
            body_contents.append(FlexText(text="📝 現場概況", weight="bold", size="sm", margin="md"))
            body_contents.append(FlexText(text=summary[:300], wrap=True, size="sm", color="#374151"))

        # 稽核表分類項目（新格式）
        categories = scene.get("categories")
        if isinstance(categories, dict):
            body_contents.extend(_flex_category_contents(categories))
        else:
            # 舊格式相容
            risks = [str(r).strip() for r in (scene.get("risks") or []) if str(r).strip()]
            if risks:
                body_contents.append(FlexSeparator(margin="md"))
                body_contents.append(FlexText(text="⚠️ 風險", weight="bold", size="sm", color="#dc2626", margin="md"))
                for r in risks[:5]:
                    body_contents.append(FlexText(text=f"• {r[:120]}", wrap=True, size="sm", color="#374151"))

        next_action = str(scene.get("next_action") or "").strip()
        if next_action:
            body_contents.append(FlexSeparator(margin="md"))
            body_contents.append(FlexText(text="⚖️ 改善建議", weight="bold", size="sm", margin="md"))
            body_contents.append(FlexText(text=next_action[:300], wrap=True, size="sm", color="#374151"))

        answer = str(scene.get("answer") or "").strip()
        if answer and answer != summary:
            body_contents.append(FlexSeparator(margin="md"))
            body_contents.append(FlexText(text="💡 法規說明", weight="bold", size="sm", margin="md"))
            body_contents.append(FlexText(text=answer[:300], wrap=True, size="sm", color="#374151"))

        refs = scene.get("regulation_refs") or []
        if refs:
            flat = [
                str(r.get("cite", r.get("source", str(r)))) if isinstance(r, dict) else str(r)
                for r in refs[:3]
            ]
            body_contents.append(FlexSeparator(margin="md"))
            body_contents.append(FlexText(text="📚 " + "；".join(flat)[:200], wrap=True, size="xs", color="#6b7280", margin="md"))

        # 保底：LLM 回傳純文字（JSON 解析失敗時）
        raw_text = str(scene.get("raw_text") or "").strip()
        if raw_text and len(body_contents) <= 1:
            body_contents.append(FlexSeparator(margin="md"))
            body_contents.append(FlexText(text=raw_text[:400], wrap=True, size="sm", color="#374151", margin="md"))

        footer_text = f"紀錄 #{audit_id}" if audit_id else datetime.now().strftime("%Y-%m-%d %H:%M")

        bubble = FlexBubble(
            header=FlexBox(
                layout="vertical",
                contents=[FlexText(
                    text=f"🔍 AI 職安稽核結果　{risk_emoji} {risk_level}風險",
                    weight="bold", color="#ffffff", size="md",
                )],
                background_color=header_color,
                padding_all="12px",
            ),
            body=FlexBox(
                layout="vertical",
                contents=body_contents,
                padding_all="16px",
                spacing="none",
            ),
            footer=FlexBox(
                layout="vertical",
                contents=[FlexText(text=footer_text, size="xs", color="#9ca3af", align="end")],
                padding_all="8px",
            ),
        )
        return FlexMessage(alt_text="AI 職安稽核結果", contents=bubble)
    except Exception as e:
        print(f"[flex] 建立 FlexMessage 失敗：{e}")
        return None


# =============================================================================
# LINE 回應邏輯
# =============================================================================

_HELP_TEXT = (
    "【AI 職安稽核系統】可用功能：\n\n"
    "📷 照片稽核 — 傳送照片，AI 自動偵測違規\n"
    "📋 稽核紀錄 — 按時間查詢歷史紀錄，輸入「統計」看整體數據\n"
    "📚 法規查詢 — 輸入問題查詢相關法規條文\n"
    "📊 案例學習 — 分析常見違規與每日稽核趨勢\n"
    "🗂 知識庫   — 管理法規知識庫（新增/刪除/匯入）\n\n"
    "也可輸入「風險分析」或「標誌辨識」切換分析模式。\n"
    "直接傳送照片會以當前模式分析。"
)

_HELP_MODES = [
    ("📷", "照片稽核", "傳送照片，AI 自動偵測違規"),
    ("📋", "稽核紀錄", "按時間查詢，輸入「統計」看整體數據"),
    ("📚", "法規查詢", "輸入問題查詢相關法規條文"),
    ("📊", "案例學習", "分析常見違規與每日稽核趨勢"),
    ("🗂", "知識庫",   "管理法規知識庫（新增/刪除/匯入）"),
]


def build_flex_help() -> FlexMessage | None:
    if not _FLEX_OK:
        return None
    bubbles = []
    for icon, name, desc in _HELP_MODES:
        bubbles.append(FlexBubble(
            size="kilo",
            header=FlexBox(
                layout="vertical",
                contents=[FlexText(text=f"{icon} {name}", weight="bold", color="#ffffff", size="md")],
                background_color="#6366f1",
                padding_all="10px",
            ),
            body=FlexBox(
                layout="vertical",
                contents=[FlexText(text=desc, wrap=True, size="sm", color="#374151")],
                padding_all="12px",
            ),
            footer=FlexBox(
                layout="vertical",
                contents=[FlexButton(
                    action=MessageAction(label=f"切換到{name}", text=name),
                    style="primary",
                    color="#6366f1",
                    height="sm",
                )],
                padding_all="8px",
            ),
        ))
    return FlexMessage(alt_text="AI 職安稽核系統功能選單", contents=FlexCarousel(contents=bubbles))


def reply_help():
    return build_flex_help() or _HELP_TEXT


def reply_regulation_query(text: str, state: dict, page: int = 0) -> str:
    if not KNOWLEDGE_BASE_AVAILABLE:
        return f"法規查詢模組尚未安裝，無法查詢「{text}」。"

    if text != state.get("page_query") or state.get("page_mode") != MODE_REGULATION_QUERY:
        snippets = retrieve_regulations(text, regulations_dir=REGULATIONS_DIR, top_k=25)
        state["page_items"] = snippets
        state["page_query"] = text
        state["page_mode"] = MODE_REGULATION_QUERY
        state["page_index"] = 0
        page = 0
    else:
        snippets = state["page_items"]

    if not snippets:
        return f"找不到與「{text}」相關的法規內容。"

    total = len(snippets)
    total_pages = (total + PAGE_SIZE - 1) // PAGE_SIZE
    start = page * PAGE_SIZE
    end = min(start + PAGE_SIZE, total)

    blocks: list[str] = []
    for i, s in enumerate(snippets[start:end], start + 1):
        body = s.text.strip()
        if len(body) > 700:
            body = body[:700].rstrip() + "…"
        header = f"{i}. {s.source}"
        if getattr(s, "article", ""):
            header += f"｜{s.article}"
        blocks.append(f"{header}\n{body}")

    result = f"查詢：{text}（第 {page + 1}/{total_pages} 頁，共 {total} 條）\n\n"
    result += "\n\n---\n\n".join(blocks)

    nav: list[str] = []
    if end < total:
        nav.append("輸入「下一頁」查看更多")
    if page > 0:
        nav.append("輸入「上一頁」返回上頁")
    if nav:
        result += "\n\n" + " | ".join(nav)

    return result[:4900]


_AUDIT_STAT_COMMANDS = {"統計", "stats", "總覽"}
_AUDIT_TODAY_COMMANDS = {"今日", "today"}
_AUDIT_WEEK_COMMANDS  = {"本週", "本周", "week"}


def _audit_date_summary(label: str, date_prefix: str) -> str:
    normalized = normalize_audit_time_query(date_prefix)
    if not normalized:
        return f"找不到「{label}」的稽核紀錄。"
    with db_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                SELECT COUNT(*) AS total, SUM(detection_count) AS detections,
                       MIN(created_at) AS earliest, MAX(created_at) AS latest
                FROM audit_records WHERE created_at_digits LIKE %s
                """,
                (f"%{normalized}%",),
            )
            row = cur.fetchone() or {}
    if not row.get("total"):
        return f"找不到「{label}」的稽核紀錄。"
    return "\n".join([
        f"📋 {label}稽核摘要",
        f"稽核筆數：{row['total']} 筆",
        f"偵測項目：{row['detections'] or 0} 項",
        f"最早：{str(row['earliest'] or '')[:16] or '-'}",
        f"最晚：{str(row['latest'] or '')[:16] or '-'}",
        "\n輸入時間查詢逐筆紀錄，或輸入「統計」查整體數據。",
    ])


def _audit_overall_stats() -> str:
    seven_days_ago = (datetime.now() - timedelta(days=7)).strftime("%Y-%m-%d")
    with db_conn() as conn:
        with conn.cursor() as cur:
            cur.execute("SELECT COUNT(*) AS total, SUM(detection_count) AS total_det FROM audit_records")
            ov = cur.fetchone() or {}
            cur.execute(
                """
                SELECT SUBSTR(created_at, 1, 10) AS day,
                       COUNT(*) AS cnt, SUM(detection_count) AS detections
                FROM audit_records WHERE SUBSTR(created_at, 1, 10) >= %s
                GROUP BY day ORDER BY day DESC
                """,
                (seven_days_ago,),
            )
            recent = cur.fetchall()
    if not ov.get("total"):
        return "尚無稽核紀錄。"
    lines = [f"📊 稽核統計\n總筆數：{ov['total']} 筆 / 偵測總項：{ov['total_det'] or 0} 項"]
    if recent:
        lines.append("\n【近七日每日】")
        for row in recent:
            lines.append(f"  {row['day']}：{row['cnt']} 筆 / {row['detections'] or 0} 項")
    lines.append("\n輸入時間查詢逐筆紀錄（如 2026-05-01）。")
    return "\n".join(lines)[:4900]


def format_audit_records_reply(time_query: str, source_key: str, state: dict, page: int = 0) -> str:
    cleaned = str(time_query or "").strip()
    if not cleaned:
        return "請輸入要查詢的時間，例如 2026-03-31、2026-03-31 14 或 1430。\n或輸入「統計」查看整體數據。"
    if not DB_AVAILABLE:
        return "資料庫尚未連線，請在 .env 設定 DATABASE_URL 後重啟。"

    if cleaned in _AUDIT_STAT_COMMANDS:
        return _audit_overall_stats()
    if cleaned in _AUDIT_TODAY_COMMANDS:
        return _audit_date_summary("今日", datetime.now().strftime("%Y-%m-%d"))
    if cleaned in _AUDIT_WEEK_COMMANDS:
        return _audit_date_summary("本週", (datetime.now() - timedelta(days=6)).strftime("%Y-%m-%d"))

    state["page_query"] = cleaned
    state["page_mode"] = MODE_AUDIT_RECORD

    offset = page * PAGE_SIZE
    records = find_audit_records(cleaned, source_key=source_key, limit=PAGE_SIZE + 1, offset=offset)
    has_more = len(records) > PAGE_SIZE
    state["has_more_pages"] = has_more
    records = records[:PAGE_SIZE]

    if not records and page == 0:
        return f"找不到符合「{cleaned}」的稽核紀錄。"
    if not records:
        return "已是最後一頁，沒有更多紀錄。"

    lines = [f"稽核紀錄（第 {page + 1} 頁）："]
    for idx, rec in enumerate(records, offset + 1):
        summary = str(rec.get("summary") or "").strip().replace("\n", " / ") or "無摘要"
        attachments: list[str] = []
        if rec.get("has_original_image"):
            attachments.append("原圖")
        if rec.get("has_annotated_image"):
            attachments.append("標記圖")
        if rec.get("has_reference_image"):
            attachments.append("比對圖")
        att = f" | 附件：{'/'.join(attachments)}" if attachments else ""
        lines.append(
            f"{idx}. #{rec.get('id', '')} {rec.get('created_at', '')} | "
            f"偵測 {rec.get('detection_count', 0)} 項 | "
            f"{summary[:150]}{att}"
        )

    nav: list[str] = []
    if has_more:
        nav.append("輸入「下一頁」查看更多")
    if page > 0:
        nav.append("輸入「上一頁」返回上頁")
    if nav:
        lines.append("\n" + " | ".join(nav))

    return "\n".join(lines)[:4900]


_CASE_WEEK_COMMANDS = {"本週", "本周", "week"}
_CASE_TOP_COMMANDS  = {"最常見", "top", "top5"}


_KW_EN = {
    "安全帽": "Helmet", "安全帶": "Harness", "鷹架": "Scaffold",
    "護欄": "Guardrail", "反光背心": "Reflect Vest", "安全鞋": "Safety Shoes",
    "電氣": "Electrical", "火災": "Fire", "墜落": "Fall", "感電": "Electrocution",
}


def generate_case_chart(stats: dict, host_url: str) -> str | None:
    """生成違規關鍵字 + 每日趨勢長條圖，回傳公開 URL；失敗回傳 None。"""
    if not _MATPLOTLIB_OK or not host_url or not stats:
        return None
    try:
        kw_data = [(_KW_EN.get(k, k), int(c or 0)) for k, c in stats.get("kw_counts", []) if int(c or 0) > 0][:8]
        trend_data = list(reversed(stats.get("trend", [])[:7]))

        if not kw_data and not trend_data:
            return None

        n = int(bool(kw_data)) + int(bool(trend_data))
        fig, axes = _plt.subplots(n, 1, figsize=(10, 4 * n))
        if n == 1:
            axes = [axes]

        ax_idx = 0
        if kw_data:
            ax = axes[ax_idx]; ax_idx += 1
            labels = [k for k, _ in kw_data]
            counts = [c for _, c in kw_data]
            bars = ax.barh(labels, counts, color="#ef4444")
            ax.set_xlabel("Count")
            ax.set_title("Top Violation Keywords")
            ax.bar_label(bars, padding=3, fmt="%d")
            ax.set_xlim(0, max(counts) * 1.25)
            ax.invert_yaxis()

        if trend_data:
            ax = axes[ax_idx]
            days = [row["day"][-5:] for row in trend_data]
            counts = [int(row["cnt"] or 0) for row in trend_data]
            bars2 = ax.bar(range(len(days)), counts, color="#3b82f6")
            ax.set_xticks(range(len(days)))
            ax.set_xticklabels(days, rotation=45, ha="right")
            ax.set_ylabel("Audits")
            ax.set_title("Daily Audit Trend (Last 7 Days)")
            ax.bar_label(bars2, padding=3, fmt="%d")
            ax.set_ylim(0, max(counts) * 1.25 if counts else 5)

        _plt.tight_layout(pad=2.0)
        filename = f"case_chart_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
        fig_path = OUTPUT_DIR / filename
        _plt.savefig(fig_path, dpi=130, bbox_inches="tight")
        _plt.close(fig)
        return f"{host_url}/outputs/{filename}"
    except Exception as exc:
        print(f"[chart] 圖表生成失敗：{exc}")
        return None


def build_flex_case_study(stats: dict, period: str) -> FlexMessage | None:
    if not _FLEX_OK or not stats:
        return None
    try:
        ov = stats.get("overview") or {}
        total = ov.get("total") or 0
        if not total:
            return None

        body_contents: list = [
            FlexText(text=f"稽核次數：{total} 次", size="sm", color="#374151"),
            FlexText(text=f"平均偵測：{ov.get('avg_det') or 0} 項　最高：{ov.get('max_det') or 0} 項",
                     size="sm", color="#374151"),
        ]

        kw = [(k, int(c or 0)) for k, c in stats.get("kw_counts", []) if int(c or 0) > 0][:5]
        if kw:
            body_contents.append(FlexSeparator(margin="md"))
            body_contents.append(FlexText(text="⚠️ 常見違規", weight="bold", size="sm", margin="md"))
            max_c = max(c for _, c in kw) or 1
            for rank, (keyword, count) in enumerate(kw, 1):
                bar = "█" * round(count / max_c * 8) + "░" * (8 - round(count / max_c * 8))
                body_contents.append(
                    FlexText(text=f"{rank}. {keyword} {bar} {count}次",
                             size="sm", color="#374151", wrap=False)
                )

        trend = stats.get("trend", [])[:5]
        if trend:
            body_contents.append(FlexSeparator(margin="md"))
            body_contents.append(FlexText(text="📅 近期趨勢", weight="bold", size="sm", margin="md"))
            for row in reversed(trend):
                body_contents.append(
                    FlexText(text=f"{row['day']}　{row['cnt']}次 / {row.get('detections') or 0}項",
                             size="sm", color="#374151")
                )

        bubble = FlexBubble(
            header=FlexBox(
                layout="vertical",
                contents=[FlexText(text=f"📊 案例分析｜{period}", weight="bold", color="#ffffff", size="md")],
                background_color="#0f766e",
                padding_all="12px",
            ),
            body=FlexBox(layout="vertical", contents=body_contents, padding_all="16px", spacing="sm"),
            footer=FlexBox(
                layout="vertical",
                contents=[FlexText(text="輸入「本週」/「本月」/「最常見」切換", size="xs", color="#9ca3af", align="center")],
                padding_all="8px",
            ),
        )
        return FlexMessage(alt_text=f"案例分析｜{period}", contents=bubble)
    except Exception as e:
        print(f"[flex_case] 建立失敗：{e}")
        return None


def reply_case_study(text: str, stats: dict | None = None) -> str:
    if not DB_AVAILABLE:
        return "案例學習需要資料庫連線，請在 .env 設定 DATABASE_URL 後重啟。"

    t = text.strip()
    days = 7 if t in _CASE_WEEK_COMMANDS else 30
    stats = stats if stats is not None else get_case_study_stats(days)
    if not stats:
        return "目前尚無稽核案例紀錄。"

    ov = stats["overview"]
    total = ov.get("total") or 0
    if not total:
        return "目前尚無稽核案例紀錄。"

    period = "本週（近 7 日）" if t in _CASE_WEEK_COMMANDS else "本月（近 30 日）"
    lines = [f"📊 案例分析｜{period}"]
    lines.append(
        f"稽核次數：{total} 次 | 平均偵測：{ov.get('avg_det') or 0} 項 | "
        f"最高偵測：{ov.get('max_det') or 0} 項 | 偵測總項：{ov.get('total_det') or 0} 項"
    )

    # 違規關鍵字排行
    kw = [(k, c) for k, c in stats.get("kw_counts", []) if c > 0]
    if t in _CASE_TOP_COMMANDS or kw:
        lines.append("\n【常見違規關鍵字】")
        if kw:
            for rank, (keyword, count) in enumerate(kw[:7], 1):
                bar = "█" * min(count, 10)
                lines.append(f"  {rank}. {keyword}：{count} 次 {bar}")
        else:
            lines.append("  （無命中關鍵字）")

    # 每日趨勢
    trend = stats.get("trend", [])
    if trend and t not in _CASE_TOP_COMMANDS:
        lines.append("\n【每日稽核趨勢】")
        for row in trend[:7]:
            lines.append(f"  {row['day']}：{row['cnt']} 次 / {row['detections'] or 0} 項")

    lines.append("\n輸入「本週」/「本月」切換期間，「最常見」只看違規排行。")
    return "\n".join(lines)[:4900]


_KB_LIST_COMMANDS = {"列出", "來源", "來源列表", "list"}
_KB_IMPORT_COMMANDS = {"匯入", "import"}


def _kb_delete_source(source_name: str) -> str:
    if not DB_AVAILABLE:
        return "資料庫未連線，無法刪除。"
    try:
        with db_conn() as conn:
            with conn.cursor() as cur:
                cur.execute("DELETE FROM regulations WHERE source = %s", (source_name,))
                deleted = cur.rowcount
            conn.commit()
        return f"已刪除「{source_name}」（{deleted} 段）。" if deleted else f"找不到來源「{source_name}」。"
    except Exception as e:
        return f"刪除失敗：{e}"


def reply_knowledge_base(text: str, state: dict) -> str:
    t = text.strip()

    if not t or t.lower() in _KB_LIST_COMMANDS:
        sources = list_regulation_sources()
        if not sources:
            return (
                "知識庫目前無資料。\n"
                "可用指令：\n"
                "「新增」— 傳送檔案新增法規\n"
                "「匯入」— 從 regulations/ 資料夾匯入 .md/.txt"
            )
        # 分頁（重新查詢時重置）
        if state.get("page_mode") != MODE_KNOWLEDGE_BASE or state.get("page_query") != "__list__":
            state["page_items"] = sources
            state["page_query"] = "__list__"
            state["page_mode"] = MODE_KNOWLEDGE_BASE
            state["page_index"] = 0
        page = state.get("page_index", 0)
        total = len(sources)
        total_pages = (total + PAGE_SIZE - 1) // PAGE_SIZE
        start = page * PAGE_SIZE
        end = min(start + PAGE_SIZE, total)
        lines = [f"📚 知識庫來源（第 {page + 1}/{total_pages} 頁，共 {total} 個）："]
        for i, row in enumerate(sources[start:end], start + 1):
            lines.append(f"{i}. {row['source']}（{row['chunks']} 段）")
        nav: list[str] = []
        if end < total:
            nav.append("輸入「下一頁」查看更多")
        if page > 0:
            nav.append("輸入「上一頁」返回上頁")
        lines.append("\n「新增」傳檔 | 「刪除 [來源]」刪除 | 「匯入」從資料夾匯入")
        if nav:
            lines.append(" | ".join(nav))
        return "\n".join(lines)[:4900]

    if t.lower() in _KB_IMPORT_COMMANDS:
        if not KNOWLEDGE_BASE_AVAILABLE:
            return "知識庫模組尚未安裝。"
        try:
            from knowledge_base import import_regulations_to_db
            count = import_regulations_to_db(REGULATIONS_DIR)
            return f"匯入完成，共新增 {count} 段法規。" if count else "regulations/ 資料夾無可匯入的檔案。"
        except Exception as e:
            return f"匯入失敗：{e}"

    if t.startswith("刪除 ") or t.startswith("刪除　"):
        source_name = t[3:].strip()
        return _kb_delete_source(source_name) if source_name else "請在「刪除」後輸入來源名稱。"

    if t == "新增":
        if not DB_AVAILABLE:
            return "資料庫未連線，無法新增。"
        state["kb_pending"] = {"action": "await_file"}
        exts = "PDF / Word / Excel / PowerPoint / TXT / Markdown / CSV / HTML"
        return f"請傳送法規文件，支援格式：\n{exts}\n\n檔案名稱將作為知識庫來源名稱。\n輸入「取消」可中止。"

    return "請輸入有效指令：「新增」、「刪除 [來源]」、「匯入」、「列出」。"


# =============================================================================
# Flask 路由（僅 LINE webhook 與靜態輸出）
# =============================================================================


@app.route("/outputs/<path:filename>")
def serve_outputs(filename):
    return send_from_directory(OUTPUT_DIR, filename)


@app.route("/audit/<int:audit_id>/annotated")
def serve_audit_annotated(audit_id: int):
    if not DB_AVAILABLE:
        return jsonify({"error": "資料庫未連線"}), 503
    with db_conn() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT annotated_image_data FROM audit_records WHERE id = %s",
                (audit_id,),
            )
            row = cur.fetchone()
    if not row or not row["annotated_image_data"]:
        return jsonify({"error": "找不到標記圖片"}), 404
    return send_file(
        io.BytesIO(bytes(row["annotated_image_data"])),
        mimetype="image/jpeg",
    )


@app.route("/export")
def export_records():
    """GET /export?month=2026-05  或  ?date=2026-05-11  下載 Excel 稽核報表。"""
    if not DB_AVAILABLE:
        return jsonify({"error": "資料庫未連線，請設定 DATABASE_URL"}), 503

    try:
        import openpyxl
        from openpyxl.styles import Alignment, Font, PatternFill
    except ImportError:
        return jsonify({"error": "請安裝 openpyxl：pip install openpyxl"}), 500

    month = request.args.get("month", "").strip()
    date_q = request.args.get("date", "").strip()
    query_str = month or date_q or ""
    limit = min(int(request.args.get("limit", 500)), 2000)

    records = find_audit_records(query_str, source_key=None, limit=limit, offset=0)

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "稽核紀錄"

    headers = ["ID", "時間", "來源", "稽核模式", "偵測數量", "分析摘要"]
    col_widths = [8, 20, 10, 14, 10, 60]
    hdr_fill = PatternFill(start_color="1A1A2E", end_color="1A1A2E", fill_type="solid")
    hdr_font = Font(color="FFFFFF", bold=True)

    for col_idx, (header, width) in enumerate(zip(headers, col_widths), 1):
        cell = ws.cell(row=1, column=col_idx, value=header)
        cell.fill = hdr_fill
        cell.font = hdr_font
        cell.alignment = Alignment(horizontal="center", vertical="center")
        ws.column_dimensions[cell.column_letter].width = width
    ws.row_dimensions[1].height = 24

    for rec in records:
        ws.append([
            rec.get("id"),
            rec.get("created_at", ""),
            rec.get("source", ""),
            rec.get("analysis_mode", ""),
            rec.get("detection_count", 0),
            (rec.get("summary") or "").replace("\n", " / ")[:500],
        ])

    ws.freeze_panes = "A2"

    ws2 = wb.create_sheet("摘要統計")
    ws2.column_dimensions["A"].width = 14
    ws2.column_dimensions["B"].width = 30
    for row in [
        ["匯出時間", datetime.now().strftime("%Y-%m-%d %H:%M:%S")],
        ["查詢條件", query_str or "（全部）"],
        ["總筆數", len(records)],
        ["偵測項目總計", sum(r.get("detection_count", 0) for r in records)],
    ]:
        ws2.append(row)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    tag = query_str.replace("-", "") if query_str else "all"
    filename = f"audit_{tag}_{datetime.now().strftime('%Y%m%d')}.xlsx"
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        download_name=filename,
        as_attachment=True,
    )


def _handle_line_webhook_post():
    signature = request.headers.get("X-Line-Signature", "")
    body = request.get_data().decode("utf-8")

    try:
        handler.handle(body, signature)
    except InvalidSignatureError:
        print(f"❌ Signature 驗證失敗 | secret={LINE_CHANNEL_SECRET[:8]}... sig={signature[:20]}")
        return "OK", 200
    except Exception as e:
        print(f"❌ 其他錯誤: {e!r}")
        return "OK", 200

    return "OK", 200


@app.route("/callback", methods=["POST"])
def callback():
    return _handle_line_webhook_post()


@app.route("/line/webhook", methods=["GET", "POST"])
def line_webhook():
    if request.method == "GET":
        return jsonify({"ok": True, "message": "LINE webhook endpoint is reachable"})
    return _handle_line_webhook_post()


# =============================================================================
# LINE 事件處理器
# =============================================================================

_QR_IMAGE_MODES = (("📷 照片稽核", "照片稽核"), ("⚠️ 風險分析", "風險分析"), ("🔍 標誌辨識", "標誌辨識"), ("📚 法規查詢", "法規查詢"))
_QR_MAIN_MENU   = (("📷 稽核", "照片稽核"), ("📚 法規", "法規查詢"), ("📊 案例", "案例學習"), ("📋 紀錄", "稽核紀錄"), ("🗂 知識庫", "知識庫"))
_QR_CASE        = (("本週", "本週"), ("本月", "案例學習"), ("最常見", "最常見"), ("📷 稽核", "照片稽核"))
_QR_KB          = (("新增", "新增"), ("列出", "列出"), ("匯入", "匯入"), ("📷 稽核", "照片稽核"))
_QR_AUDIT       = (("統計", "統計"), ("今日", "今日"), ("本週", "本週"), ("📷 稽核", "照片稽核"))


def _page_qr(state: dict, extra: tuple = ()) -> QuickReply | None:
    """Build pagination Quick Reply based on current page state."""
    items: list[tuple] = list(extra)
    if state.get("page_index", 0) > 0:
        items.insert(0, ("⬅ 上一頁", "上一頁"))
    has_next = state.get("has_more_pages") or (
        state.get("page_items") and
        (state["page_index"] + 1) * PAGE_SIZE < len(state["page_items"])
    )
    if has_next:
        items.insert(0, ("➡ 下一頁", "下一頁"))
    return _qr(*items) if items else None


@handler.add(MessageEvent, message=TextMessageContent)
def handle_text(event):
    text = (event.message.text or "").strip()
    source_key = get_source_key(event)
    state = get_user_state(source_key)

    # ── 幫助選單 ──────────────────────────────────────────────────────────
    if not text or text in HELP_COMMANDS:
        msg = reply_help()
        if isinstance(msg, FlexMessage):
            _line_reply(event.reply_token, msg)
        else:
            _line_reply(event.reply_token, TextMessage(text=msg, quick_reply=_qr(*_QR_MAIN_MENU)))
        return

    # ── 知識庫「新增」等待上傳檔案中 ──────────────────────────────────────
    if state.get("kb_pending") and state["kb_pending"].get("action") == "await_file":
        if text in {"取消", "cancel"}:
            state["kb_pending"] = None
            _line_reply(event.reply_token, TextMessage(text="已取消新增。", quick_reply=_qr(*_QR_KB)))
        else:
            _line_reply(event.reply_token, TextMessage(
                text="請傳送檔案（PDF/Word/Excel 等），或輸入「取消」中止。"
            ))
        return

    # ── 切換模式 ──────────────────────────────────────────────────────────
    if text in MODE_COMMANDS:
        cmd = MODE_COMMANDS[text]
        set_user_mode(source_key, cmd["mode"], cmd.get("prompt"))
        state["page_items"] = []
        state["page_index"] = 0
        state["page_query"] = ""
        state["page_mode"] = ""
        state["kb_pending"] = None
        label = _MODE_LABELS.get(cmd["mode"], cmd["mode"])
        hint = _MODE_HINTS.get(cmd["mode"], "")
        mode_qr = {
            MODE_IMAGE_ANALYSIS:   _QR_IMAGE_MODES,
            MODE_REGULATION_QUERY: (("📷 稽核", "照片稽核"), ("📋 紀錄", "稽核紀錄"), ("📊 案例", "案例學習")),
            MODE_CASE_STUDY:       _QR_CASE,
            MODE_AUDIT_RECORD:     _QR_AUDIT,
            MODE_KNOWLEDGE_BASE:   _QR_KB,
        }.get(cmd["mode"], _QR_MAIN_MENU)
        _line_reply(event.reply_token, TextMessage(
            text=f"已切換到「{label}」模式。\n{hint}",
            quick_reply=_qr(*mode_qr),
        ))
        return

    # ── 翻頁（下一頁） ────────────────────────────────────────────────────
    if text in NEXT_PAGE_COMMANDS:
        page_mode = state.get("page_mode", "")
        if page_mode == MODE_REGULATION_QUERY:
            state["page_index"] += 1
            reply = reply_regulation_query(state["page_query"], state, state["page_index"])
            _line_reply(event.reply_token, TextMessage(
                text=reply,
                quick_reply=_page_qr(state, (("📷 稽核", "照片稽核"),)),
            ))
        elif page_mode == MODE_AUDIT_RECORD:
            state["page_index"] += 1
            reply = format_audit_records_reply(state["page_query"], source_key, state, state["page_index"])
            _line_reply(event.reply_token, TextMessage(
                text=reply,
                quick_reply=_page_qr(state, _QR_AUDIT),
            ))
        elif page_mode == MODE_KNOWLEDGE_BASE:
            state["page_index"] += 1
            reply = reply_knowledge_base("列出", state)
            _line_reply(event.reply_token, TextMessage(
                text=reply,
                quick_reply=_page_qr(state, _QR_KB),
            ))
        else:
            _line_reply(event.reply_token, TextMessage(text="目前沒有可翻頁的內容。", quick_reply=_qr(*_QR_MAIN_MENU)))
        return

    # ── 翻頁（上一頁） ────────────────────────────────────────────────────
    if text in PREV_PAGE_COMMANDS:
        page_mode = state.get("page_mode", "")
        if state.get("page_index", 0) <= 0:
            _line_reply(event.reply_token, TextMessage(text="已是第一頁。", quick_reply=_qr(*_QR_MAIN_MENU)))
            return
        state["page_index"] -= 1
        if page_mode == MODE_REGULATION_QUERY:
            reply = reply_regulation_query(state["page_query"], state, state["page_index"])
            _line_reply(event.reply_token, TextMessage(
                text=reply,
                quick_reply=_page_qr(state, (("📷 稽核", "照片稽核"),)),
            ))
        elif page_mode == MODE_AUDIT_RECORD:
            reply = format_audit_records_reply(state["page_query"], source_key, state, state["page_index"])
            _line_reply(event.reply_token, TextMessage(
                text=reply,
                quick_reply=_page_qr(state, _QR_AUDIT),
            ))
        elif page_mode == MODE_KNOWLEDGE_BASE:
            reply = reply_knowledge_base("列出", state)
            _line_reply(event.reply_token, TextMessage(
                text=reply,
                quick_reply=_page_qr(state, _QR_KB),
            ))
        else:
            _line_reply(event.reply_token, TextMessage(text="目前沒有可翻頁的內容。", quick_reply=_qr(*_QR_MAIN_MENU)))
        return

    # ── 各模式處理 ────────────────────────────────────────────────────────
    mode = state["mode"]

    if mode == MODE_REGULATION_QUERY:
        state["page_index"] = 0
        reply = reply_regulation_query(text, state, 0)
        _line_reply(event.reply_token, TextMessage(
            text=reply,
            quick_reply=_page_qr(state, (("📷 稽核", "照片稽核"), ("📋 紀錄", "稽核紀錄"))),
        ))

    elif mode == MODE_AUDIT_RECORD:
        state["page_index"] = 0
        reply = format_audit_records_reply(text, source_key, state, 0)
        _line_reply(event.reply_token, TextMessage(
            text=reply,
            quick_reply=_page_qr(state, _QR_AUDIT),
        ))

    elif mode == MODE_CASE_STUDY:
        _t = text.strip()
        _days = 7 if _t in _CASE_WEEK_COMMANDS else 30
        _period = "本週" if _t in _CASE_WEEK_COMMANDS else ("最常見" if _t in _CASE_TOP_COMMANDS else "本月")
        _stats = get_case_study_stats(_days) if DB_AVAILABLE else {}
        _host = PUBLIC_URL or f"{request.headers.get('X-Forwarded-Proto','https')}://{request.host}"
        _chart_url = generate_case_chart(_stats, _host)
        flex_card = build_flex_case_study(_stats, _period)
        case_qr = _qr(*_QR_CASE)
        if flex_card:
            if _chart_url:
                _line_reply(event.reply_token,
                            ImageMessage(original_content_url=_chart_url, preview_image_url=_chart_url),
                            flex_card)
            else:
                _line_reply(event.reply_token, flex_card)
        else:
            reply = reply_case_study(text, _stats)
            if _chart_url:
                _line_reply(event.reply_token,
                            ImageMessage(original_content_url=_chart_url, preview_image_url=_chart_url),
                            TextMessage(text=reply, quick_reply=case_qr))
            else:
                _line_reply(event.reply_token, TextMessage(text=reply, quick_reply=case_qr))

    elif mode == MODE_KNOWLEDGE_BASE:
        reply = reply_knowledge_base(text, state)
        _line_reply(event.reply_token, TextMessage(
            text=reply,
            quick_reply=_page_qr(state, _QR_KB) if state.get("page_mode") == MODE_KNOWLEDGE_BASE else _qr(*_QR_KB),
        ))

    else:
        state["prompt"] = text
        _line_reply(event.reply_token, TextMessage(
            text=f"已更新稽核提示詞：{text}\n接著傳送照片即可進行分析。",
            quick_reply=_qr(*_QR_IMAGE_MODES),
        ))


def _process_image_async(
    user_id: str,
    source_key: str,
    message_id: str,
    prompt: str,
    host_url: str,
) -> None:
    try:
        image_bytes = _line_get_content(message_id)
    except Exception as e:
        print(f"下載失敗: {e}")
        _line_push(user_id, TextMessage(text="❌ 圖片下載失敗。"))
        return

    try:
        result = analyze_image_bytes(
            image_bytes,
            prompt=prompt,
            use_yolo=True,
            filename_stem=f"line_{message_id}",
        )
    except Exception as e:
        print(f"分析失敗: {e}")
        _line_push(user_id, TextMessage(text="❌ AI 分析出錯。"))
        return

    try:
        audit_record = save_audit_record(
            source="line",
            source_key=source_key,
            prompt=prompt,
            original_image_data=image_bytes,
            result=result,
            analysis_mode="line image",
        )
    except Exception as e:
        print(f"儲存紀錄失敗: {e}")
        # DB 寫入失敗不應讓分析結果整個消失：退回不含紀錄 ID 的乾淨 record，
        # 讓使用者至少能收到這次的分析結果。
        audit_record = {"id": None, "annotated_image_data": b""}

    audit_id = audit_record.get("id")

    try:
        analysis_qr = _qr(*_QR_IMAGE_MODES)
        # 一律回傳純文字，不用 Flex 卡片：Flex 是圖形化氣泡卡，使用者在 LINE 裡
        # 沒辦法整段選取複製；純文字才能讓使用者複製貼上（例如轉貼給其他人或存檔）。
        reply_text = build_line_analysis_reply(result, prompt)
        if audit_id:
            reply_text += f"\n\n📋 紀錄 ID：#{audit_id}"
        result_message = TextMessage(text=reply_text, quick_reply=analysis_qr)
    except Exception as e:
        print(f"結果訊息組裝失敗: {e}")
        _line_push(user_id, TextMessage(text="❌ 分析結果格式異常，無法顯示完整報告，請稍後再試或聯絡管理員。"))
        return

    try:
        if audit_id and host_url and audit_record.get("annotated_image_data"):
            image_url = f"{host_url}/audit/{audit_id}/annotated"
            _line_push(user_id, ImageMessage(
                original_content_url=image_url,
                preview_image_url=image_url,
            ))
        _line_push(user_id, result_message)
    except Exception as e:
        print(f"發送失敗: {e}")


@handler.add(MessageEvent, message=ImageMessageContent)
def handle_image(event):
    source_key = get_source_key(event)
    state = get_user_state(source_key)
    mode = state["mode"]
    user_id = event.source.user_id

    if mode in _MODE_IMAGE_REJECTS:
        _line_reply(event.reply_token, TextMessage(text=_MODE_IMAGE_REJECTS[mode]))
        return

    prompt = state.get("prompt") or DEFAULT_ANALYZE_PROMPT
    message_id = event.message.id

    host_url = PUBLIC_URL
    if not host_url:
        protocol = request.headers.get("X-Forwarded-Proto", "https")
        host_url = f"{protocol}://{request.host}"

    _line_reply(event.reply_token, TextMessage(text="🔎 已收到照片，正在進行 AI 職安稽核分析，請稍候..."))

    threading.Thread(
        target=_process_image_async,
        args=(user_id, source_key, message_id, prompt, host_url),
        daemon=True,
    ).start()


def _process_kb_file_async(user_id: str, message_id: str, filename: str) -> None:
    try:
        file_bytes = _line_get_content(message_id)
    except Exception as e:
        _line_push(user_id, TextMessage(text=f"❌ 檔案下載失敗：{e}"))
        return

    try:
        text = extract_text_from_file(file_bytes, filename)
    except RuntimeError as e:
        _line_push(user_id, TextMessage(text=f"❌ {e}"))
        return
    except Exception as e:
        _line_push(user_id, TextMessage(text=f"❌ 檔案解析失敗：{e}"))
        return

    text = text.strip()
    if not text:
        _line_push(user_id, TextMessage(text="❌ 無法從檔案中擷取文字，請確認檔案內容。"))
        return

    try:
        from knowledge_base import add_regulation_text
        count = add_regulation_text(filename, text)
        _line_push(user_id, TextMessage(
            text=f"✅ 已將「{filename}」新增至知識庫\n段落數：{count} 段 / 字元數：{len(text):,} 字"
        ))
    except Exception as e:
        _line_push(user_id, TextMessage(text=f"❌ 知識庫儲存失敗：{e}"))


@handler.add(MessageEvent, message=FileMessageContent)
def handle_file(event):
    source_key = get_source_key(event)
    state = get_user_state(source_key)
    user_id = event.source.user_id

    if not (state.get("kb_pending") and state["kb_pending"].get("action") == "await_file"):
        _line_reply(event.reply_token, TextMessage(
            text="若要將檔案新增至知識庫，請先切換至「知識庫」模式並輸入「新增」。"
        ))
        return

    filename = event.message.file_name or "unknown_file"
    state["kb_pending"] = None

    _line_reply(event.reply_token, TextMessage(text=f"⏳ 正在解析「{filename}」，請稍候..."))

    threading.Thread(
        target=_process_kb_file_async,
        args=(user_id, event.message.id, filename),
        daemon=True,
    ).start()


# =============================================================================
# HTTPS / CLI 啟動
# =============================================================================

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="AI 職安稽核系統")
    parser.add_argument("--host", default="0.0.0.0")
    parser.add_argument("--port", type=int, default=5443)
    parser.add_argument("--http", action="store_true", help="使用 HTTP 而非 HTTPS")
    parser.add_argument("--debug", action="store_true")
    tunnel_group = parser.add_mutually_exclusive_group()
    tunnel_group.add_argument("--tunnel", action="store_true", default=True,
                              help="自動啟動 Cloudflare Tunnel（預設開啟）")
    tunnel_group.add_argument("--no-tunnel", dest="tunnel", action="store_false",
                              help="不啟動 Cloudflare Tunnel")
    return parser.parse_args()


def get_local_ip() -> str:
    sock = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
    try:
        sock.connect(("8.8.8.8", 80))
        return sock.getsockname()[0]
    except OSError:
        return "127.0.0.1"
    finally:
        sock.close()


def get_ssl_context(use_http: bool):
    if use_http:
        return None
    if not CERT_FILE.exists() or not KEY_FILE.exists():
        raise FileNotFoundError(f"找不到憑證：{CERT_FILE} 或 {KEY_FILE}")
    return str(CERT_FILE), str(KEY_FILE)


def print_startup_urls(port: int, scheme: str) -> None:
    local_ip = get_local_ip()
    print(f"本機 URL : {scheme}://localhost:{port}")
    print(f"區網 URL : {scheme}://{local_ip}:{port}")
    print(f"LINE Webhook : {scheme}://{local_ip}:{port}/line/webhook")
    print(f"資料庫   : {'已連線' if DB_AVAILABLE else '未連線（僅本地模式）'}")
    print(f"Channel Secret (前8碼): {LINE_CHANNEL_SECRET[:8]}…（請與 LINE Developer Console 核對）")


# =============================================================================
# Cloudflare Tunnel
# =============================================================================

def start_cloudflare_tunnel(local_url: str) -> None:
    global PUBLIC_URL, _cloudflared_proc

    if not Path(CLOUDFLARED_PATH).exists():
        print(f"[tunnel] 找不到 cloudflared：{CLOUDFLARED_PATH}")
        print("[tunnel] 請在 .env 設定 CLOUDFLARED_PATH 或確認安裝路徑。")
        return

    def _run() -> None:
        global PUBLIC_URL, _cloudflared_proc
        cmd = [CLOUDFLARED_PATH, "tunnel", "--url", local_url, "--no-tls-verify"]
        print(f"[tunnel] 啟動中：{' '.join(cmd)}")
        try:
            proc = subprocess.Popen(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
                encoding="utf-8",
                errors="replace",
            )
            _cloudflared_proc = proc
            url_re = re.compile(r"https://[a-zA-Z0-9\-]+\.trycloudflare\.com")
            for line in proc.stdout or []:
                line = line.rstrip()
                if line:
                    print(f"[tunnel] {line}")
                match = url_re.search(line)
                if match and not PUBLIC_URL:
                    PUBLIC_URL = match.group(0)
                    sep = "=" * 62
                    print(f"\n{sep}")
                    print(f"  Cloudflare 公開 URL : {PUBLIC_URL}")
                    print(f"  LINE Webhook        : {PUBLIC_URL}/line/webhook")
                    print(f"{sep}\n")
            proc.wait()
            print("[tunnel] cloudflared 已結束。")
        except Exception as exc:
            print(f"[tunnel] 啟動失敗：{exc}")

    threading.Thread(target=_run, daemon=True, name="cloudflared").start()


def stop_cloudflare_tunnel() -> None:
    global _cloudflared_proc
    if _cloudflared_proc and _cloudflared_proc.poll() is None:
        _cloudflared_proc.terminate()
        _cloudflared_proc = None


init_audit_db()

try:
    from knowledge_base import init_regulations_table
    init_regulations_table()
except Exception as _kb_err:
    print(f"[KB] regulations 表初始化失敗: {_kb_err}")


if __name__ == "__main__":
    args = parse_args()
    scheme = "http" if args.http else "https"
    try:
        ssl_context = get_ssl_context(args.http)
    except FileNotFoundError as e:
        print(f"警告：{e}，改用 HTTP 啟動。")
        ssl_context = None
        scheme = "http"
    print_startup_urls(args.port, scheme)
    if args.tunnel:
        local_url = f"{scheme}://127.0.0.1:{args.port}"
        start_cloudflare_tunnel(local_url)
    try:
        app.run(host=args.host, port=args.port, debug=args.debug, ssl_context=ssl_context)
    finally:
        stop_cloudflare_tunnel()
